#!/usr/bin/env python3
"""Generate a LaTeX PDF report and a PowerPoint deck from REPORT.md."""
from __future__ import annotations

import re
import subprocess
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
REPORT_MD = ROOT / "REPORT.md"
REPORTS_DIR = ROOT / "reports"
REPORTS_DIR.mkdir(exist_ok=True)

TEX_PATH = REPORTS_DIR / "smallgameagent_report.tex"
PDF_PATH = REPORTS_DIR / "smallgameagent_report.pdf"
PPTX_PATH = REPORTS_DIR / "smallgameagent_report.pptx"


# ---------------------------------------------------------------------------
# PowerPoint theme constants
# ---------------------------------------------------------------------------
_C_PRIMARY = RGBColor(0x1A, 0x23, 0x7E)       # deep indigo
_C_SECONDARY = RGBColor(0x00, 0x96, 0xC7)     # teal
_C_ACCENT = RGBColor(0xFF, 0x6B, 0x35)        # orange
_C_LIGHT = RGBColor(0xF3, 0xF6, 0xF9)         # off-white
_C_DARK = RGBColor(0x1E, 0x29, 0x33)          # near-black
_C_MUTED = RGBColor(0x6C, 0x75, 0x7D)         # gray
_C_SUCCESS = RGBColor(0x28, 0xA7, 0x45)       # green
_C_CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)       # white
_C_WARN = RGBColor(0xF4, 0xB4, 0x00)          # amber

_W = Inches(13.333)
_H = Inches(7.5)


_LATEX_SPECIAL = {
    "\\": r"\textbackslash{}",
    "{": r"\{",
    "}": r"\}",
    "$": r"\$",
    "&": r"\&",
    "%": r"\%",
    "#": r"\#",
    "_": r"\_",
    "^": r"\^{}",
    "~": r"\textasciitilde{}",
}


def esc(s: str) -> str:
    """Escape LaTeX special chars in a single pass."""
    return re.sub(r"[\\{}$&%#_^~]", lambda m: _LATEX_SPECIAL[m.group(0)], s)


def inline_fmt(s: str) -> str:
    """Convert markdown inline formatting to LaTeX without double-escaping commands."""
    placeholders: dict[str, Any] = {}
    counter = 0

    def stash(kind: str, content: Any) -> str:
        nonlocal counter
        key = f"__{kind}_{counter:04d}__"
        counter += 1
        placeholders[key] = content
        return key

    s = re.sub(r"`([^`]+)`", lambda m: stash("CODE", m.group(1)), s)
    s = re.sub(r"\*\*([^*]+)\*\*", lambda m: stash("BOLD", m.group(1)), s)
    s = re.sub(r"\*([^*]+)\*", lambda m: stash("ITALIC", m.group(1)), s)
    s = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", lambda m: stash("LINK", (m.group(1), m.group(2))), s)

    # Now safe to escape remaining plain text.
    s = esc(s)

    # Restore placeholders with properly escaped content.
    for key, val in placeholders.items():
        if key.startswith("__CODE_"):
            replacement = f"\\texttt{{{esc(val)}}}"
        elif key.startswith("__BOLD_"):
            replacement = f"\\textbf{{{esc(val)}}}"
        elif key.startswith("__ITALIC_"):
            replacement = f"\\textit{{{esc(val)}}}"
        else:  # LINK
            text, url = val
            replacement = f"\\href{{{esc(url)}}}{{{esc(text)}}}"
        s = s.replace(key, replacement)
    return s


def render_table(rows: list[str]) -> str:
    # rows like "| a | b |"
    cells = [[c.strip() for c in r.split("|")[1:-1]] for r in rows]
    if len(cells) < 3:
        return ""
    header = cells[0]
    data = cells[2:]
    n = len(header)
    cols = "|" + "|".join(["l"] * n) + "|"
    lines = ["\\begin{table}[htbp]", "\\centering", "\\small", f"\\begin{{tabular}}{{{cols}}}", "\\hline"]
    lines.append(" & ".join(inline_fmt(c) for c in header) + " \\\\")
    lines.append("\\hline")
    for row in data:
        if len(row) < n:
            row = row + [""] * (n - len(row))
        lines.append(" & ".join(inline_fmt(c) for c in row[:n]) + " \\\\")
        lines.append("\\hline")
    lines.extend(["\\end{tabular}", "\\end{table}"])
    return "\n".join(lines)


_CIRCLED_DIGITS = {chr(0x2460 + i): f"({i+1})" for i in range(20)}


def md_to_latex(md: str) -> str:
    md = "".join(_CIRCLED_DIGITS.get(ch, ch) for ch in md)
    lines = md.splitlines()
    out: list[str] = []
    in_list = False
    in_table = False
    table_rows: list[str] = []

    def close() -> None:
        nonlocal in_list, in_table, table_rows
        if in_list:
            out.append("\\end{itemize}")
            in_list = False
        if in_table:
            out.append(render_table(table_rows))
            table_rows = []
            in_table = False

    for raw in lines:
        stripped = raw.strip()
        # headings
        if stripped.startswith("# "):
            close()
            out.append(f"\\section{{{inline_fmt(stripped[2:])}}}")
            continue
        if stripped.startswith("## "):
            close()
            out.append(f"\\subsection{{{inline_fmt(stripped[3:])}}}")
            continue
        if stripped.startswith("### "):
            close()
            out.append(f"\\subsubsection{{{inline_fmt(stripped[4:])}}}")
            continue
        # table
        if stripped.startswith("|"):
            if not in_table:
                close()
                in_table = True
                table_rows = []
            table_rows.append(stripped)
            continue
        else:
            if in_table:
                out.append(render_table(table_rows))
                table_rows = []
                in_table = False
        # empty
        if not stripped:
            if in_list:
                out.append("\\end{itemize}")
                in_list = False
            out.append("")
            continue
        # bullet
        if stripped.startswith(("- ", "* ")):
            if not in_list:
                out.append("\\begin{itemize}")
                in_list = True
            out.append("  \\item " + inline_fmt(stripped[2:]))
            continue
        # normal paragraph
        if in_list:
            out.append("\\end{itemize}")
            in_list = False
        out.append(inline_fmt(stripped) + "\\\\")
    close()
    return "\n".join(out)


def write_tex(body: str) -> None:
    tex = r"""\documentclass[12pt,a4paper]{ctexart}
\usepackage[margin=2.5cm]{geometry}
\usepackage{graphicx}
\usepackage{booktabs}
\usepackage{longtable}
\usepackage{hyperref}
\usepackage{xcolor}
\usepackage{listings}
\usepackage{setspace}
\onehalfspacing

\hypersetup{
  colorlinks=true,
  linkcolor=blue,
  urlcolor=blue,
  citecolor=blue
}

\title{\textbf{smallgameagent 实验报告} \\ \large 基于 LLM/VLM 与规则的小游戏 Agent}
\author{smallgameagent 项目组}
\date{\today}

\begin{document}
\maketitle
\tableofcontents
\newpage
""" + body + "\n\\end{document}\n"
    TEX_PATH.write_text(tex, encoding="utf-8")


def compile_tex() -> None:
    for _ in range(2):
        result = subprocess.run(
            ["xelatex", "-interaction=nonstopmode", "-output-directory", str(REPORTS_DIR), str(TEX_PATH)],
            cwd=str(REPORTS_DIR),
            capture_output=True,
            text=True,
        )
        if result.returncode != 0:
            (REPORTS_DIR / "xelatex.log").write_text(result.stdout + "\n" + result.stderr, encoding="utf-8")
            print("xelatex failed; see reports/xelatex.log", file=sys.stderr)
            raise SystemExit(result.returncode)
    print(f"PDF -> {PDF_PATH}")


# ---------------------------------------------------------------------------
# PowerPoint helpers
# ---------------------------------------------------------------------------
def _set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.width = Pt(0)


def _set_text_style(
    paragraph,
    text: str,
    font_size: Pt,
    bold: bool = False,
    color: RGBColor = _C_DARK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = "Microsoft YaHei",
) -> None:
    paragraph.text = text
    paragraph.font.size = font_size
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.font.name = font_name
    paragraph.alignment = align


def _add_left_bar(slide, color: RGBColor = _C_ACCENT, width: float = 0.18) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.05), Inches(width), Inches(6.45))
    _set_fill(bar, color)


def _add_header_bar(slide, title: str) -> None:
    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), _W, Inches(1.05))
    _set_fill(bar, _C_PRIMARY)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.6))
    _set_text_style(tb.text_frame.paragraphs[0], title, Pt(30), bold=True, color=_C_LIGHT)


def _add_footer(slide, page: int, total: int) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.15), _W, Inches(0.02))
    _set_fill(line, _C_SECONDARY)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(4), Inches(0.25))
    _set_text_style(tb.text_frame.paragraphs[0], "smallgameagent · LLM/VLM + 规则驱动", Pt(10), color=_C_MUTED)
    tb2 = slide.shapes.add_textbox(Inches(11.8), Inches(7.2), Inches(1), Inches(0.25))
    _set_text_style(tb2.text_frame.paragraphs[0], f"{page}/{total}", Pt(10), color=_C_MUTED, align=PP_ALIGN.RIGHT)


def _add_content_bg(slide) -> None:
    """Subtle off-white background panel for content slides."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(1.18), Inches(12.63), Inches(5.85))
    _set_fill(bg, _C_CARD_BG)
    bg.line.color.rgb = RGBColor(0xE1, 0xE5, 0xE9)
    bg.line.width = Pt(1)


def _card(slide, left: float, top: float, width: float, height: float, title: str, bullets: list[str], accent: RGBColor = _C_SECONDARY) -> None:
    # Subtle shadow shape (offset slightly, muted gray)
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + 0.03), Inches(top + 0.03), Inches(width), Inches(height))
    _set_fill(shadow, RGBColor(0xD0, 0xD5, 0xDA))
    shadow.fill.transparency = 0.7

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    _set_fill(shape, _C_CARD_BG)
    shape.line.color.rgb = RGBColor(0xDD, 0xE2, 0xE6)
    shape.line.width = Pt(1)
    # Title strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.45))
    _set_fill(strip, accent)
    strip.line.width = Pt(0)
    tb_title = slide.shapes.add_textbox(Inches(left + 0.08), Inches(top + 0.08), Inches(width - 0.16), Inches(0.35))
    _set_text_style(tb_title.text_frame.paragraphs[0], title, Pt(14), bold=True, color=_C_LIGHT)
    # Body
    tb = slide.shapes.add_textbox(Inches(left + 0.12), Inches(top + 0.55), Inches(width - 0.24), Inches(height - 0.65))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_text_style(p, b, Pt(12), color=_C_DARK)
        p.space_after = Pt(5)


def _add_table_slide(prs, title: str, headers: list[str], rows: list[list[str]], row_heights: list[float] | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, title)
    _add_left_bar(slide, _C_SECONDARY)
    _add_content_bg(slide)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    left = Inches(0.65)
    top = Inches(1.35)
    width = Inches(12.1)
    height = Inches(5.5)
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = _C_PRIMARY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.bold = True
            paragraph.font.color.rgb = _C_LIGHT
            paragraph.font.name = "Microsoft YaHei"
            paragraph.alignment = PP_ALIGN.CENTER
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _C_LIGHT
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.name = "Microsoft YaHei"
                paragraph.alignment = PP_ALIGN.CENTER
    return slide


def _add_bullet_slide(prs, title: str, bullets: list[str]) -> Any:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, title)
    _add_left_bar(slide, _C_ACCENT)
    _add_content_bg(slide)
    content = slide.shapes.add_textbox(Inches(0.85), Inches(1.45), Inches(11.8), Inches(5.5))
    tf = content.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_text_style(p, b, Pt(18), color=_C_DARK)
        p.space_after = Pt(12)
    return slide


def _add_title_slide(prs, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Deep indigo base with subtle bottom gradient
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), _W, _H)
    _set_fill(bg, _C_PRIMARY)

    # Top diagonal accent (translucent teal strip)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(-1), Inches(6), Inches(10))
    _set_fill(strip, _C_SECONDARY)
    strip.fill.fore_color.brightness = 0.25

    # Decorative circles
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(0.9), Inches(2.0), Inches(2.0))
    _set_fill(circle1, _C_ACCENT)
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.6), Inches(2.2), Inches(1.1), Inches(1.1))
    _set_fill(circle2, _C_LIGHT)

    # Bottom wave bar
    bar1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(-0.5), Inches(5.7), Inches(14.5), Inches(2.1))
    _set_fill(bar1, _C_SECONDARY)
    bar2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(-0.5), Inches(6.1), Inches(14.5), Inches(1.7))
    _set_fill(bar2, RGBColor(0x00, 0x7A, 0xA3))

    # Thin orange line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(5.55), Inches(4.33), Inches(0.06))
    _set_fill(line, _C_ACCENT)

    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.5))
    _set_text_style(tb.text_frame.paragraphs[0], title, Pt(54), bold=True, color=_C_LIGHT, align=PP_ALIGN.CENTER)

    # Subtitle
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.65), Inches(12.3), Inches(1.1))
    _set_text_style(tb2.text_frame.paragraphs[0], subtitle, Pt(24), color=_C_LIGHT, align=PP_ALIGN.CENTER)

    # Tagline at bottom
    tb3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.6))
    _set_text_style(tb3.text_frame.paragraphs[0], "分层多 Agent · 在线规则进化 · 批量数据管线", Pt(16), bold=True, color=_C_LIGHT, align=PP_ALIGN.CENTER)


def _add_two_column_slide(prs, title: str, left_title: str, left_bullets: list[str], right_title: str, right_bullets: list[str]) -> Any:
    """A slide with two cards side by side."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, title)
    _add_left_bar(slide, _C_SECONDARY)
    _add_content_bg(slide)
    _card(slide, 0.65, 1.45, 6.0, 5.4, left_title, left_bullets, _C_SECONDARY)
    _card(slide, 6.95, 1.45, 6.0, 5.4, right_title, right_bullets, _C_ACCENT)
    return slide


def _add_rule_wiring_slide(prs) -> None:
    """Visual slide showing RuleParameters shared between L2 and L0."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "规则在线更新：L2 的参数能真正改变 L0")
    _add_left_bar(slide, _C_SUCCESS)
    _add_content_bg(slide)

    # L2 box
    l2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.6), Inches(4.0), Inches(1.4))
    _set_fill(l2, _C_PRIMARY)
    _set_text_style(l2.text_frame.paragraphs[0], "L2 云端 API\n检测触发条件，输出 param 更新", Pt(14), bold=True, color=_C_LIGHT, align=PP_ALIGN.CENTER)

    # Arrow down
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(2.6), Inches(3.1), Inches(0.8), Inches(0.7))
    _set_fill(arrow, _C_MUTED)

    # Shared parameter store
    store = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.9), Inches(4.0), Inches(1.0))
    _set_fill(store, _C_SECONDARY)
    _set_text_style(store.text_frame.paragraphs[0], "共享 RuleParameters\ncoin_save_buffer / stuck_escape_threshold / ...", Pt(13), bold=True, color=_C_LIGHT, align=PP_ALIGN.CENTER)

    # Arrow down
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(2.6), Inches(4.95), Inches(0.8), Inches(0.7))
    _set_fill(arrow2, _C_MUTED)

    # L0 box
    l0 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.75), Inches(4.0), Inches(1.2))
    _set_fill(l0, _C_ACCENT)
    _set_text_style(l0.text_frame.paragraphs[0], "L0 规则引擎\n每步读取参数，执行 move / tap", Pt(14), bold=True, color=_C_LIGHT, align=PP_ALIGN.CENTER)

    # Side note card
    _card(slide, 6.5, 1.6, 6.3, 5.35, "本次修复的关键点", [
        "之前 L2 的 param 更新只停留在 HierarchicalPlanner 内部",
        "RuleEngine 用硬编码默认值，导致更新无效",
        "现在 HybridAgent 创建共享 RuleParameters，同时注入 L2 和 L0",
        "L2 更新 → 共享参数 → L0 即时生效，形成完整闭环",
    ], _C_SUCCESS)


def _add_trigger_design_slide(prs) -> None:
    """Slide explaining how rule updates are triggered."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "规则更新怎么触发？阈值 + 两层安全门")
    _add_left_bar(slide, _C_WARN)
    _add_content_bg(slide)

    _card(slide, 0.65, 1.45, 6.0, 2.6, "触发条件（满足任一即上报）", [
        "composite 连续 N 步低于阈值（默认 0.15）",
        "stall 停滞计数超过阈值（默认 5 步）",
        "L0 与 L2 决策连续 K 步冲突",
        "世界模型检测到空间/时间一致性违例",
    ], _C_ACCENT)

    _card(slide, 6.95, 1.45, 6.0, 2.6, "L1 本地 VLM 的作用", [
        "每 N 步或卡住时看截图",
        "输出「场景类型 / 关键目标 / 障碍 / UI」",
        "把视觉上下文拼进 L2 prompt，辅助判断",
        "不直接改规则，只提供改规则的证据",
    ], _C_SECONDARY)

    # Bottom flow
    _card(slide, 0.65, 4.25, 12.3, 2.5, "更新决策流", [
        "① 触发器收集信号 → ② L1 补充视觉证据 → ③ L2 输出结构化 JSON（update_type / target / payload / confidence）",
        "④ Applier 安全门：allowlist + 置信度 ≥ 0.9 + patch ≤ 2000 字符 + search 唯一匹配 + 自动备份",
        "⑤ 通过则写入 runtime_rules.json；未通过进待审队列，由人工或更高阈值模型二次确认",
        "⑥ RuleEngine 下一步读取新参数，行为立即改变；重启后配置仍然有效",
    ], _C_PRIMARY)


def _add_code_file_update_slide(prs) -> None:
    """Slide showing code-file rule update experiment."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "代码文件级更新：L2 直接调引擎旋钮")
    _add_left_bar(slide, _C_SUCCESS)
    _add_content_bg(slide)

    _card(slide, 0.65, 1.45, 6.0, 2.6, "安全门设计", [
        "白名单：只能改 configs/runtime_rules.json",
        "置信度 ≥ 0.9，patch ≤ 2000 字符",
        "search 块必须唯一匹配，否则进待审队列",
        "修改前自动备份，保留最近 3 份",
    ], _C_SUCCESS)

    _card(slide, 6.95, 1.45, 6.0, 2.6, "离线实验结果", [
        "游戏：SSD_00461P01，回放 30 步",
        "mock L2 以 0.95 置信度触发更新",
        "stuck_escape_threshold 5 → 3",
        "第 6 步起规则引擎读取到新阈值",
        "实验结束后配置文件自动恢复",
    ], _C_ACCENT)

    # Provider table
    n_rows = 5
    n_cols = 4
    left = Inches(0.65)
    top = Inches(4.35)
    width = Inches(12.1)
    height = Inches(1.5)
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    headers = ["Provider", "模型", "延迟", "结果"]
    rows = [
        ["qwen", "qwen3.7-max", "7.83s", "✅ 应用成功"],
        ["kimi", "kimi-k2.7-code", "3.32s", "✅ 应用成功"],
        ["xiaomi", "mimo-v2.5", "6.18s", "✅ 应用成功"],
        ["opencodego", "deepseek-v4-flash", "9.08s", "✅ 应用成功"],
    ]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = _C_PRIMARY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = _C_LIGHT
            paragraph.font.name = "Microsoft YaHei"
            paragraph.alignment = PP_ALIGN.CENTER
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _C_LIGHT
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Microsoft YaHei"
                paragraph.alignment = PP_ALIGN.CENTER

    insight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(5.95), Inches(12.1), Inches(0.55))
    _set_fill(insight, _C_LIGHT)
    _set_text_style(insight.text_frame.paragraphs[0],
        "意义：云端模型不仅能改内存参数，还能持久化地调整规则引擎的「旋钮」，且 qwen/kimi/xiaomi/opencodego 四家都能正确生成可应用的 JSON patch。",
        Pt(12), color=_C_DARK, align=PP_ALIGN.CENTER)


def _add_agent_comm_slide(prs) -> None:
    """Agent communication and memory strategies slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "Agent 通信与记忆：从单兵到协作")
    _add_left_bar(slide, _C_SECONDARY)
    _add_content_bg(slide)

    _card(slide, 0.55, 1.45, 4.0, 2.3, "记忆三层", [
        "Episodic：逐步 state / action / reward",
        "Semantic：游戏类型、目标语义向量",
        "Procedural：成功策略、失败模式",
    ], _C_PRIMARY)

    _card(slide, 4.85, 1.45, 4.0, 2.3, "AgentBus 消息", [
        "OBSERVATION / STATE_UPDATE",
        "DECISION_PROPOSAL / CRITIQUE",
        "RULE_UPDATE / MEMORY_WRITE",
    ], _C_SECONDARY)

    _card(slide, 9.15, 1.45, 4.0, 2.3, "协作流程", [
        "Observer 采集 probe 状态",
        "DecisionAnalyst 生成候选动作",
        "Critic 评估一致性风险",
        "MemoryCurator 归档有用模式",
    ], _C_ACCENT)

    # Bottom insight
    insight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(4.0), Inches(12.5), Inches(1.1))
    _set_fill(insight, _C_LIGHT)
    _set_text_style(insight.text_frame.paragraphs[0],
        "核心思路：不要把所有决策压给一个模型。让云端 API 做长程规划，本地 VLM 做画面理解，规则引擎做零延迟执行，记忆层负责跨 episode 复用。",
        Pt(14), color=_C_DARK, align=PP_ALIGN.CENTER)


def _add_provider_slide(prs) -> None:
    """Updated multi-provider slide."""
    _add_table_slide(prs, "云端多 Provider 实测状态",
        ["Provider", "文本模型", "视觉模型", "状态"],
        [
            ["Kimi", "kimi-k2.7-code / k2.5", "kimi-k2.6 / k2.5", "✅ 可用"],
            ["MiMo (xiaomi)", "mimo-v2.5", "mimo-v2.5", "✅ 可用"],
            ["OpenCodeGo", "deepseek-v4-flash", "mimo-v2.5", "✅ 可用"],
            ["Qwen", "qwen3.7-max", "qwen3.7-max", "文本 ✅ / 视觉 ⚠️"],
            ["DeepSeek", "deepseek-chat", "deepseek-chat", "❌ 余额不足"],
        ])


def _add_metric_card(slide, left: float, top: float, width: float, height: float, number: str, label: str, accent: RGBColor = _C_PRIMARY) -> None:
    """A large-number metric card for impact slides."""
    # Shadow
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + 0.02), Inches(top + 0.02), Inches(width), Inches(height))
    _set_fill(shadow, RGBColor(0xD0, 0xD5, 0xDA))
    shadow.fill.transparency = 0.65
    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    _set_fill(card, _C_CARD_BG)
    card.line.color.rgb = accent
    card.line.width = Pt(2)
    # Accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.12))
    _set_fill(bar, accent)
    bar.line.width = Pt(0)
    # Number
    tb_num = slide.shapes.add_textbox(Inches(left), Inches(top + 0.25), Inches(width), Inches(height * 0.55))
    _set_text_style(tb_num.text_frame.paragraphs[0], number, Pt(36), bold=True, color=accent, align=PP_ALIGN.CENTER)
    # Label
    tb_label = slide.shapes.add_textbox(Inches(left + 0.05), Inches(top + height * 0.55), Inches(width - 0.1), Inches(height * 0.35))
    _set_text_style(tb_label.text_frame.paragraphs[0], label, Pt(12), color=_C_MUTED, align=PP_ALIGN.CENTER)


def _add_metrics_slide(prs) -> None:
    """Slide with 4-5 big numbers summarizing the project."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "项目速览：我们做到了什么")
    _add_left_bar(slide, _C_SUCCESS)
    _add_content_bg(slide)

    metrics = [
        ("22", "款游戏已可驱动", _C_PRIMARY),
        ("15,083", "条 VLM 训练样本", _C_SECONDARY),
        ("0.251", "rule 基线 mean composite", _C_SUCCESS),
        ("4", "家云端 provider 可用", _C_ACCENT),
        ("3", "层架构（L0/L1/L2）", _C_PRIMARY),
    ]

    # Layout: top row 3 cards, bottom row 2 cards centered
    positions = [
        (0.9, 1.6, 3.7, 2.0),
        (4.85, 1.6, 3.7, 2.0),
        (8.8, 1.6, 3.7, 2.0),
        (2.85, 3.95, 3.7, 2.0),
        (6.8, 3.95, 3.7, 2.0),
    ]
    for (left, top, w, h), (num, label, color) in zip(positions, metrics):
        _add_metric_card(slide, left, top, w, h, num, label, color)

    # Bottom insight
    insight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(6.15), Inches(11.6), Inches(0.7))
    _set_fill(insight, _C_LIGHT)
    _set_text_style(insight.text_frame.paragraphs[0],
        "核心思路：用规则打底保证零延迟，云端 API 和本地 VLM 只在必要时出手修正。",
        Pt(14), color=_C_DARK, align=PP_ALIGN.CENTER)


def _add_humanized_intro_slide(prs) -> None:
    """A more narrative 'why this matters' slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "为什么这件事值得做？")
    _add_left_bar(slide, _C_SECONDARY)
    _add_content_bg(slide)

    _card(slide, 0.65, 1.45, 6.0, 2.5, "小游戏可玩广告自动化", [
        "每天有大量 playable ad 需要测试和调优",
        "手动玩、手动评效率低，难以规模化",
        "不同游戏后端差异大，'一刀切'策略走不通",
    ], _C_PRIMARY)

    _card(slide, 6.95, 1.45, 6.0, 2.5, "我们的解法", [
        "把'快执行'和'慢思考'拆开：规则负责毫秒级动作",
        "云端大模型负责长程规划和规则进化",
        "本地小 VLM 作为'眼睛'，给云端提供画面证据",
    ], _C_SECONDARY)

    # Arrow between cards
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.05), Inches(2.35), Inches(1.5), Inches(0.7))
    _set_fill(arrow, _C_ACCENT)

    # Bottom quote-like box
    quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(4.25), Inches(12.3), Inches(1.9))
    _set_fill(quote, _C_LIGHT)
    tf = quote.text_frame
    tf.word_wrap = True
    points = [
        "目标：让 Agent 能像人一样理解游戏画面、记住有效策略、在卡住时自我调整",
        "同时保持每步接近零延迟，不拖累 playable ad 的实时交互",
    ]
    for i, txt in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_text_style(p, txt, Pt(16), color=_C_DARK, align=PP_ALIGN.CENTER)
        p.space_after = Pt(8)


def _add_section_slide(prs, number: str, title: str, subtitle: str = "") -> None:
    """A full-bleed section divider with a dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), _W, _H)
    _set_fill(bg, _C_PRIMARY)

    # Diagonal accent
    diag = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(-1), Inches(6.5), Inches(9.5))
    _set_fill(diag, _C_SECONDARY)
    diag.fill.fore_color.brightness = 0.2

    # Large number watermark
    num_tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(3), Inches(2))
    _set_text_style(num_tb.text_frame.paragraphs[0], number, Pt(120), bold=True, color=RGBColor(0x33, 0x3D, 0xA0))

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(3.1), Inches(12.1), Inches(1.2))
    _set_text_style(tb.text_frame.paragraphs[0], title, Pt(44), bold=True, color=_C_LIGHT)
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(4.4), Inches(12.1), Inches(0.8))
        _set_text_style(tb2.text_frame.paragraphs[0], subtitle, Pt(22), color=_C_LIGHT)

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.4 if not subtitle else 5.1), Inches(2), Inches(0.08))
    _set_fill(accent, _C_ACCENT)


def _add_architecture_slide(prs) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "三层架构：规则打底，云端/本地 VLM 做上层更新")
    _add_left_bar(slide, _C_PRIMARY)
    _add_content_bg(slide)

    # Draw three layered boxes
    layers = [
        ("L2 战略层 · 云端多模态 API", "kimi-k2.7-code / mimo-v2.5 / DeepSeek / Qwen\n长程规划 + 规则更新触发", _C_PRIMARY, 1.4),
        ("L1 战术层 · 本地小 VLM", "qwen3.5-4b / gemma4-e4b（4-bit + KV-cache）\n离线标注 + 视觉上下文", _C_SECONDARY, 3.1),
        ("L0 执行层 · 规则引擎", "tap-guide / tap-only / 障碍学习\n零延迟执行", _C_ACCENT, 4.8),
    ]
    for title, body, color, top in layers:
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(top), Inches(6.3), Inches(1.35))
        _set_fill(rect, color)
        tb = slide.shapes.add_textbox(Inches(3.7), Inches(top + 0.12), Inches(5.9), Inches(0.45))
        _set_text_style(tb.text_frame.paragraphs[0], title, Pt(16), bold=True, color=_C_LIGHT, align=PP_ALIGN.CENTER)
        tb2 = slide.shapes.add_textbox(Inches(3.7), Inches(top + 0.55), Inches(5.9), Inches(0.7))
        _set_text_style(tb2.text_frame.paragraphs[0], body, Pt(12), color=_C_LIGHT, align=PP_ALIGN.CENTER)

    # Connecting arrows between layers
    for top in [2.75, 4.45]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.25), Inches(top), Inches(0.6), Inches(0.35))
        _set_fill(arrow, _C_MUTED)

    # Side cards
    _card(slide, 0.4, 1.5, 2.7, 2.2, "触发条件", [
        "composite 连续低于阈值",
        "stall 超过 5 步",
        "L0/L2 决策冲突",
        "世界模型 stale 命中",
    ], _C_MUTED)
    _card(slide, 10.2, 1.5, 2.7, 2.2, "更新产物", [
        "参数更新",
        "phase contract",
        "strategy memory",
        "（可选）代码片段",
    ], _C_SUCCESS)

    # Bus / memory at bottom
    bus = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(6.35), Inches(8.3), Inches(0.55))
    _set_fill(bus, _C_DARK)
    tb = slide.shapes.add_textbox(Inches(2.7), Inches(6.42), Inches(7.9), Inches(0.45))
    _set_text_style(tb.text_frame.paragraphs[0], "AgentBus · Observer → StateMapper → DecisionAnalyst → Verifier → Critic → MemoryCurator", Pt(12), color=_C_LIGHT, align=PP_ALIGN.CENTER)


def _add_roadmap_slide(prs) -> None:
    """Timeline/roadmap slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "技术路线图：从能跑到好用")
    _add_left_bar(slide, _C_SECONDARY)
    _add_content_bg(slide)

    milestones = [
        ("已完成", "规则引擎 + 自动校准\n22 游戏可驱动", _C_SUCCESS, 0.9),
        ("已完成", "multi-bus-memory\n记忆读回 0.300", _C_SUCCESS, 3.4),
        ("已完成", "L2 code-file 更新\nqwen/kimi/xiaomi/opencodego 验证", _C_SUCCESS, 5.9),
        ("进行中", "本地 VLM 上下文\nQLoRA 微调准备", _C_ACCENT, 8.4),
        ("下一步", "真实游戏在线触发\n规则更新 + A/B 回归", _C_SECONDARY, 10.9),
    ]

    # Timeline bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(3.5), Inches(12.0), Inches(0.08))
    _set_fill(bar, _C_MUTED)

    for label, desc, color, left in milestones:
        # Dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(3.35), Inches(0.35), Inches(0.35))
        _set_fill(dot, color)
        # Label above
        label_tb = slide.shapes.add_textbox(Inches(left - 0.55), Inches(2.3), Inches(1.8), Inches(0.5))
        _set_text_style(label_tb.text_frame.paragraphs[0], label, Pt(13), bold=True, color=color, align=PP_ALIGN.CENTER)
        # Desc below
        desc_tb = slide.shapes.add_textbox(Inches(left - 0.8), Inches(3.8), Inches(2.4), Inches(1.3))
        _set_text_style(desc_tb.text_frame.paragraphs[0], desc, Pt(12), color=_C_DARK, align=PP_ALIGN.CENTER)

    # Bottom note
    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(5.55), Inches(12.0), Inches(0.9))
    _set_fill(note, _C_LIGHT)
    _set_text_style(note.text_frame.paragraphs[0],
        "核心原则：把贵且慢的云端调用摊到多步，本地 VLM 只提供「证据」，规则引擎负责零延迟执行。",
        Pt(14), color=_C_DARK, align=PP_ALIGN.CENTER)


def _add_game_matrix_slide(prs) -> None:
    """Game coverage matrix slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "游戏覆盖矩阵：22 款游戏已可驱动")
    _add_left_bar(slide, _C_SECONDARY)
    _add_content_bg(slide)

    categories = [
        ("A 类 joystick", "7 款", "00440 / 00461 / 00483 / 00496 / 00517 / 00522 / 00736", _C_PRIMARY),
        ("B 类 tap-only", "15 款", "00219 / 00332 / 00342 / 00382 / 00394 / 00427 / 00434 / 00475 / 00482 / 00526 / 00532 / 00594 / 00669 / 00733 / 00742", _C_SECONDARY),
    ]

    for idx, (cat, count, games, color) in enumerate(categories):
        top = 1.5 + idx * 2.4
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(top), Inches(12.1), Inches(2.1))
        _set_fill(rect, _C_CARD_BG)
        rect.line.color.rgb = color
        rect.line.width = Pt(2)
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(top + 0.2), Inches(1.8), Inches(0.45))
        _set_fill(tag, color)
        _set_text_style(tag.text_frame.paragraphs[0], f"{cat} · {count}", Pt(13), bold=True, color=_C_LIGHT, align=PP_ALIGN.CENTER)
        tb = slide.shapes.add_textbox(Inches(0.85), Inches(top + 0.75), Inches(11.7), Inches(1.2))
        _set_text_style(tb.text_frame.paragraphs[0], games, Pt(14), color=_C_DARK)


def _add_local_vlm_slide(prs) -> None:
    """Local VLM slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "本地 VLM：8 GB 显存上的画面理解层")
    _add_left_bar(slide, _C_SECONDARY)
    _add_content_bg(slide)

    _card(slide, 0.65, 1.45, 6.0, 2.8, "当前能跑通什么？", [
        "4-bit 量化加载 Qwen3.5-4B/9B、Gemma-4-E4B",
        "KV-cache 量化进一步压低显存",
        "RTX 5060 8GB 上 Qwen3.5-4B 约 7s/帧，Gemma-4-E4B 约 7.8s/帧",
        "Gemma-4-E4B 视觉摘要比 Qwen3.5-4B 更稳定",
    ], _C_SECONDARY)

    _card(slide, 6.95, 1.45, 6.0, 2.8, "视觉上下文有没有用？", [
        "Qwen3.5-4B：text-only 5/9，with-visual 3/9",
        "Gemma-4-E4B：text-only 3/9，with-visual 5/9",
        "说明模型选对很重要，视觉摘要能帮云端纠偏",
        "下一步：QLoRA 微调，输出结构化视觉上下文",
    ], _C_ACCENT)

    insight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(4.45), Inches(12.3), Inches(1.8))
    _set_fill(insight, _C_LIGHT)
    tf = insight.text_frame
    tf.word_wrap = True
    points = [
        "结论：本地 VLM 可以跑，但默认模型还不足以稳定提升云端策略质量。",
        "方向：用 15,083 条 processed-runs 样本做 QLoRA 微调，让模型学会输出「箭头方向 + 关键目标 + 障碍物」的结构化上下文。",
        "部署：微调后通过 llama.cpp server 常驻，L1 每 N 步或卡住时提供视觉证据。",
    ]
    for i, txt in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_text_style(p, txt, Pt(13), color=_C_DARK)
        p.space_after = Pt(6)


def _add_key_findings_slide(prs) -> None:
    """Key findings summary slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "关键发现：什么有效，什么还需打磨")
    _add_left_bar(slide, _C_SUCCESS)
    _add_content_bg(slide)

    findings = [
        ("✅ 有效", "rule 在 6 游戏 representative subset 上 mean composite=0.251，仍是当前最稳短程基线。", _C_SUCCESS),
        ("✅ 有效", "multi-bus-memory 在 tap-only 游戏上接近 rule（0.296 vs 0.300），记忆读回对无 joystick 场景有价值。", _C_SUCCESS),
        ("✅ 有效", "L2 code-file 更新在 qwen/kimi/xiaomi/opencodego 均成功，云端模型确实能改持久化规则旋钮。", _C_SUCCESS),
        ("✅ 已修复", "00483 multi-bus activity=0：根因是 strategy_memory 在线自强化，session 隔离后 multi-bus-memory 0.200、activity 0.333。", _C_SUCCESS),
        ("⚠️ 待优化", "multi-bus 在需要 joystick 的 A 组明显弱于 rule（mean 0.110），需要更强的 driver/Verifier 循环。", _C_WARN),
        ("📌 方向", "把云端放在 L2 做规划与规则更新，L0 规则负责执行，L1 VLM 只做视觉证据。", _C_SECONDARY),
    ]

    for idx, (badge, text, color) in enumerate(findings):
        top = 1.45 + idx * 0.85
        badge_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(top), Inches(1.3), Inches(0.42))
        _set_fill(badge_shape, color)
        _set_text_style(badge_shape.text_frame.paragraphs[0], badge, Pt(12), bold=True, color=_C_LIGHT, align=PP_ALIGN.CENTER)
        tb = slide.shapes.add_textbox(Inches(2.2), Inches(top), Inches(10.5), Inches(0.9))
        _set_text_style(tb.text_frame.paragraphs[0], text, Pt(15), color=_C_DARK)


def _add_representative_results_slide(prs) -> None:
    """Representative subset online results slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "Representative Subset：6 游戏在线跑测")
    _add_left_bar(slide, _C_SECONDARY)
    _add_content_bg(slide)

    _add_table_slide_raw(slide, 0.65, 1.45, 12.3, 4.8,
        ["游戏", "类型", "rule", "multi-bus", "multi-bus-memory"],
        [
            ["SSD_00461P01 塔防", "A", "0.149", "0.050", "0.044"],
            ["SSD_00483P01 吸沙抽水", "A", "0.244", "0.103", "0.200"],
            ["SSD_00522P02 地下炸矿", "A", "0.215", "0.178", "0.178"],
            ["SSD_00382P01 低坑杀鲨鱼", "B", "0.300", "—", "0.288"],
            ["SSD_00594P02 破石收水", "B", "0.300", "—", "0.300"],
            ["SSD_00742P01 加油小镇", "B", "0.300", "—", "0.300"],
        ])

    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(6.35), Inches(12.3), Inches(0.7))
    _set_fill(note, _C_LIGHT)
    _set_text_style(note.text_frame.paragraphs[0],
        "15 runs 全部成功。rule mean composite=0.251 仍是最稳基线；multi-bus-memory 在 tap-only 游戏接近 rule（0.296 vs 0.300）。session 隔离让 00483 从 activity=0 恢复到 0.333。",
        Pt(13), color=_C_DARK, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Internal table helper (positionable)
# ---------------------------------------------------------------------------
def _add_table_slide_raw(slide, left: float, top: float, width: float, height: float,
                         headers: list[str], rows: list[list[str]]) -> None:
    """Add a simple table at a fixed position on an existing slide."""
    n_cols = len(headers)
    n_rows = len(rows)
    table = slide.shapes.add_table(n_rows + 1, n_cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    # Header
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = _C_PRIMARY
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.font.bold = True
        paragraph.font.color.rgb = _C_LIGHT
        paragraph.alignment = PP_ALIGN.CENTER
    # Rows
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = _C_DARK
            paragraph.alignment = PP_ALIGN.CENTER


def _add_rule_update_flow_slide(prs) -> None:
    """Visual flow diagram of the online rule-update loop."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "规则在线更新闭环：从触发到回滚")
    _add_left_bar(slide, _C_SUCCESS)
    _add_content_bg(slide)

    # Boxes: (left, top, width, height, label, color)
    boxes = [
        (0.8, 1.6, 2.6, 1.0, "L0 Rule Engine\n零延迟执行", _C_ACCENT),
        (4.2, 1.6, 2.6, 1.0, "Trigger\ncomposite / stall / conflict", _C_WARN),
        (7.6, 1.6, 2.6, 1.0, "L2 Cloud API\n结构化 JSON 输出", _C_PRIMARY),
        (7.6, 3.4, 2.6, 1.0, "Applier\nallowlist + confidence + backup", _C_SECONDARY),
        (4.2, 3.4, 2.6, 1.0, "Watchdog\nbaseline vs trial 对比", _C_SUCCESS),
        (0.8, 3.4, 2.6, 1.0, "RuleParameters\n内存 + runtime_rules.json", _C_MUTED),
    ]

    for left, top, w, h, label, color in boxes:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
        _set_fill(box, color)
        box.line.width = Pt(0)
        tb = slide.shapes.add_textbox(Inches(left + 0.08), Inches(top + 0.15), Inches(w - 0.16), Inches(h - 0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(label.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _set_text_style(p, line, Pt(12), bold=(i == 0), color=_C_LIGHT, align=PP_ALIGN.CENTER)

    # Arrows (simplified: right arrows between top row, down arrows, left arrows)
    arrow_positions = [
        (3.4, 1.95, 0.8, 0.35, MSO_SHAPE.RIGHT_ARROW),   # L0 -> Trigger
        (6.8, 1.95, 0.8, 0.35, MSO_SHAPE.RIGHT_ARROW),   # Trigger -> L2
        (8.9, 2.6, 0.35, 0.8, MSO_SHAPE.DOWN_ARROW),     # L2 -> Applier
        (6.8, 3.75, 0.8, 0.35, MSO_SHAPE.LEFT_ARROW),    # Applier -> Watchdog
        (3.4, 3.75, 0.8, 0.35, MSO_SHAPE.LEFT_ARROW),    # Watchdog -> Params
        (1.95, 2.6, 0.35, 0.8, MSO_SHAPE.UP_ARROW),      # Params -> L0
    ]
    for left, top, w, h, shape_type in arrow_positions:
        arrow = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(w), Inches(h))
        _set_fill(arrow, _C_MUTED)

    # Bottom insight
    insight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.6))
    _set_fill(insight, _C_LIGHT)
    tf = insight.text_frame
    tf.word_wrap = True
    points = [
        "① Trigger 监控 composite / stall / conflict，满足条件时上报",
        "② L2 输出结构化 JSON（param / memory_entry / code_file），Applier 安全门过滤",
        "③ Watchdog 对比更新前后 3 步表现，变差则自动 rollback",
        "④ RuleParameters 同时被 L0 和 Trigger 读取，形成「执行 → 监控 → 更新 → 回滚」闭环",
    ]
    for i, txt in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_text_style(p, txt, Pt(12), color=_C_DARK)
        p.space_after = Pt(4)


def _add_bar_chart_slide(prs) -> None:
    """Slide with horizontal bars comparing search/plan variants."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "规划变体对比：长 horizon 收益最大")
    _add_left_bar(slide, _C_SECONDARY)
    _add_content_bg(slide)

    # Data: (label, value, color)
    data = [
        ("rule 基线", 0.194, _C_MUTED),
        ("beam_2step", 0.254, _C_WARN),
        ("hierarchical_mock_5", 0.328, _C_SECONDARY),
        ("hierarchical_short", 0.328, _C_SECONDARY),
        ("hierarchical_mock_15", 0.298, _C_SECONDARY),
        ("hierarchical_long", 0.388, _C_SUCCESS),
    ]
    max_val = 0.45
    left = 3.2
    top_start = 1.7
    bar_height = 0.55
    gap = 0.25
    max_width = 8.5

    for i, (label, val, color) in enumerate(data):
        top = top_start + i * (bar_height + gap)
        # Label
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(top + 0.08), Inches(2.6), Inches(bar_height))
        _set_text_style(tb.text_frame.paragraphs[0], label, Pt(13), color=_C_DARK, align=PP_ALIGN.RIGHT)
        # Bar background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(max_width), Inches(bar_height))
        _set_fill(bg, _C_LIGHT)
        bg.line.color.rgb = RGBColor(0xE1, 0xE5, 0xE9)
        bg.line.width = Pt(1)
        # Bar fill
        width = max(0.1, val / max_val * max_width)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(bar_height))
        _set_fill(bar, color)
        bar.line.width = Pt(0)
        # Value label
        vb = slide.shapes.add_textbox(Inches(left + width + 0.1), Inches(top + 0.08), Inches(1.0), Inches(bar_height))
        _set_text_style(vb.text_frame.paragraphs[0], f"{val:.3f}", Pt(12), bold=True, color=color)

    # Insight box
    insight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.8))
    _set_fill(insight, _C_LIGHT)
    _set_text_style(insight.text_frame.paragraphs[0],
        "SSD_00461P01 · type_match rate · hierarchical_long（8 意图）比 rule 提升 +100%，比 beam search 提升 +53%",
        Pt(13), color=_C_DARK, align=PP_ALIGN.CENTER)


def _add_harness_slide(prs) -> None:
    """Slide summarizing the fps-play-agent-harness integration."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "同学框架集成：统一实验标准")
    _add_left_bar(slide, _C_SUCCESS)
    _add_content_bg(slide)

    _card(slide, 0.65, 1.45, 6.0, 2.5, "规划器桥接（planner-http-adapter）", [
        "harness_http provider ↔ 我们的 Chat Completions API",
        "deepseek-v4-flash / mimo-v2.5 / kimi-k2.7 / qwen3.7-max",
        "kimi few-shot 可产出 schema 合规 StrategySpec（7-48s）",
        "fallback-first 模式：0ms 返回确定性策略",
    ], _C_SECONDARY)

    _card(slide, 6.95, 1.45, 6.0, 2.5, "已解决的工程问题", [
        "WSL WebGL：headful + xvfb + swiftshader 标志",
        "probe 假阳性：download/ad 按钮不再误判通关（§5.2 移植）",
        "objective unresolved：无 guide 时用 target_id",
        "代理网络：NODE_USE_ENV_PROXY=1 走 clash",
    ], _C_SUCCESS)

    # Results row
    _card(slide, 0.65, 4.15, 12.3, 2.4, "运行结果与瓶颈", [
        "tiles-survive：可玩 3-4 步，joystick 校准被游戏特殊控制卡住（同学 Codex 用 1h17m 探索）",
        "whiteout 商店：SETTLED_COMPLETE 假阳性已修复 → BLOCKED_UNKNOWN_MECHANIC（诚实信号）",
        "瓶颈：模型 schema 遵循度（fallback 兜底）+ 每游戏控制知识 + WSL 软件渲染不稳定",
        "25+25 游戏清单已筛选交付（指引箭头 + 通关画面 + 广告图标验证）",
    ], _C_ACCENT)


def _add_next_steps_slide(prs) -> None:
    """Next steps slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "下一步：把实验变成稳定系统")
    _add_left_bar(slide, _C_SECONDARY)
    _add_content_bg(slide)

    bullets = [
        "在真实游戏运行中触发 code-file 更新，验证动态 stall/composite 下 L2 的决策质量",
        "把 runtime_rules.json 的 schema 写进 L2 prompt，约束可改字段与取值范围",
        "让本地 VLM 常驻，验证 L1 战术修正对 joystick 游戏的实际收益",
        "在 5090 服务器跑 QLoRA 微调，把本地 VLM 训练成专用画面理解器",
        "把离线回放接入 CI：每次规则改动自动跑 5 游戏回归",
        "探索 Agent 间更严格的仲裁：Critic 对 L0/L2 冲突做最终决策",
        " redesign strategy_memory 的 success 信号：结合位置变化、guide 触发、stall 减少",
    ]
    content = slide.shapes.add_textbox(Inches(0.85), Inches(1.55), Inches(12.0), Inches(5.5))
    tf = content.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.05), Inches(0), Inches(0.15), Inches(0.15))
        _set_fill(prefix, _C_SECONDARY)
        _set_text_style(p, f"  {b}", Pt(18), color=_C_DARK)
        p.space_after = Pt(12)


def write_pptx() -> None:
    prs = Presentation()
    prs.slide_width = _W
    prs.slide_height = _H

    # ---- Slide 1: Title ----
    _add_title_slide(prs, "smallgameagent 实验进展", "LLM/VLM + 规则驱动的小游戏 Agent\n多 Provider、在线规则更新与批量数据管线")

    # ---- Slide 2: Humanized intro ----
    _add_humanized_intro_slide(prs)

    # ---- Slide 3: Agenda ----
    _add_bullet_slide(prs, "今天聊什么？", [
        "小游戏 Agent 到底卡在哪几个地方？",
        "我们的思路：把快执行和慢思考拆开，做三层架构",
        "怎么接入多家云模型，又怎么在本地跑小 VLM",
        "规则不是死的：云端模型如何在线更新底层规则",
        "Agent 之间怎么通信、怎么记东西",
        "最近实验跑出了什么结果，数据怎么准备",
        "下一步还想试什么？",
    ])

    # ---- Slide 4: Key metrics ----
    _add_metrics_slide(prs)

    # ---- Section: Problem ----
    _add_section_slide(prs, "01", "我们在解决什么问题？")

    _add_bullet_slide(prs, "小游戏可玩广告，自动化并不 trivial", [
        "空间一致性：场景随进度变化，云端模型容易记错障碍物位置",
        "时间一致性：后续策略可能反过来改前面的行为含义，导致推进失败",
        "策略短视：有钱就升级，而不是攒够再升级，来回跑很多冤枉步",
        "运行效率：不同游戏后端数据差异大，探针和日志不能一刀切",
        "API 延迟：纯云端决策太慢，实时控制扛不住",
    ])

    # ---- Section: Architecture ----
    _add_section_slide(prs, "02", "三层架构")

    _add_architecture_slide(prs)

    _add_bullet_slide(prs, "慢思考 + 快执行，各司其职", [
        "L0 规则引擎：每步几乎不花时间，负责 move / tap 的零延迟执行",
        "L1 本地 VLM：看截图判断当前状态，每 N 步或卡住时做战术修正",
        "L2 云端 API：看 probe state 做长程规划和规则更新，每 M 步或阶段切换触发",
        "核心思路：把贵且慢的调用摊到多步，单步延迟压到接近 0",
    ])

    # ---- Section: Cloud API ----
    _add_section_slide(prs, "03", "云端多 Provider API")

    _add_provider_slide(prs)

    _add_bullet_slide(prs, "一个客户端，切换多家云模型", [
        "统一接入 OpenCodeGo / MiMo、Kimi、DeepSeek、Xiaomi、Qwen",
        ".env 集中管理 key 与 base_url，provider 用 CLOUD_PROVIDER 环境变量切换",
        "支持 KIMI_TEXT_MODEL / KIMI_VISION_MODEL 等覆盖默认模型",
        "当前实测可用：Kimi / Xiaomi / OpenCodeGo 文本+多模态；Qwen 文本可用；DeepSeek 余额不足",
        "Kimi 系列自动省略 temperature，避免代理返回 400",
    ])

    # ---- Section: Rule update ----
    _add_section_slide(prs, "04", "规则在线更新")

    _add_rule_wiring_slide(prs)

    _add_trigger_design_slide(prs)

    _add_code_file_update_slide(prs)

    _add_rule_update_flow_slide(prs)

    _add_bullet_slide(prs, "规则不是写死的，触发后才让上层改", [
        "触发器监控：composite 持续低迷 / stall 计数 / L0-L2 决策冲突 / 世界模型 stale",
        "L2 输出结构化 JSON：param、memory_entry、phase_contract、code_file",
        "默认只修改内存参数与 strategy memory，零风险、即时生效",
        "代码文件改写需 allowlist + 置信度 ≥ 0.9 + 自动备份，没通过就进待审队列",
        "目标：把「有钱就升级」的短视行为改成「攒够再升级」的长程策略",
    ])

    _add_bullet_slide(prs, "触发器改进：不再一刀切", [
        "旧问题：绝对阈值 0.15 在低基线游戏上几乎每 5 步就触发，L2 调用泛滥",
        "相对下降触发：记录本 run 峰值，下降超过 25% 才触发，避免过度敏感",
        "硬上限：单 run 最多更新 3 次，防止「越更新越差、越差越更新」",
        "cooldown 从 5 步提高到 8 步，给每次更新更长的观察窗口",
        "单元测试验证：peak 0.30 → 0.20 触发（drop 22.2%），max=2 后不再触发",
    ])

    _add_bullet_slide(prs, "Watchdog 回滚：不止看分数", [
        "旧版只比较 composite 平均值，小幅参数调整带来的 stall 增加被忽略",
        "新版三层安全网：composite 下降 / activity 下降 ≥ 0.15 / stall 增加 ≥ 2",
        "过程指标（stall/activity）恶化往往早于结果指标，提前回滚减少无效步数",
        "实验验证：stall 从 0 升到 3 时即使 composite 不变也触发回滚",
        "意义：L2 可以大胆试错，watchdog 负责兜底",
    ])

    # ---- Section: Agent communication ----
    _add_section_slide(prs, "05", "Agent 通信与记忆")

    _add_agent_comm_slide(prs)

    _add_two_column_slide(prs, "通信与记忆的设计取舍",
        "记忆读回的优势",
        [
            "00461 multi-bus-memory 0.106 → 0.300",
            "00736 multi-bus-memory 0.269 → 0.300",
            "跨 episode 复用成功模式",
            "减少重复探索带来的 stall",
        ],
        "仍需解决的问题",
        [
            "00483 / 00517 multi-bus activity=0：driver 不匹配",
            "strategy_memory 需要预热，冷启动时可能反噬",
            "多 Agent 冲突时需要更严格的仲裁机制",
        ])

    # ---- Section: Local VLM ----
    _add_section_slide(prs, "06", "本地 VLM：把 8 GB 显存用到极限")

    _add_local_vlm_slide(prs)

    # ---- Section: Results ----
    _add_section_slide(prs, "07", "实验结果")

    _add_table_slide(prs, "离线回放：rule 基线 vs 在线规则更新（5 款游戏）",
        ["游戏", "rule composite", "hierarchical(mock)", "rule updates"],
        [
            ["00461 塔防", "0.241", "0.241", "13"],
            ["00219 养牛卖奶", "0.300", "0.300", "38"],
            ["00332 圣诞薅羊毛", "0.300", "0.283", "9"],
            ["00342 建造合并", "0.300", "0.297", "35"],
            ["00382 低坑杀鲨鱼", "0.300", "0.286", "15"],
        ])

    _add_table_slide(prs, "Qwen L2 真实云端：SSD_00461P01（29 步）",
        ["指标", "数值"],
        [
            ["type_match", "10/29"],
            ["action_match", "7/29"],
            ["composite", "0.246"],
            ["mean_latency_ms", "3728.6"],
            ["rule_update_count", "4"],
        ])

    _add_key_findings_slide(prs)

    _add_representative_results_slide(prs)

    _add_bar_chart_slide(prs)

    # ---- Section: Data ----
    _add_section_slide(prs, "08", "训练数据管线")

    _add_bullet_slide(prs, "跑过的轨迹 = 可复用的训练数据", [
        "batch_runner 每步记录 state / action / keyNumbers / reason",
        "processed_runs_converter 从 22 个游戏的 processed-runs 生成 7 任务样本",
        "当前数据集：15,083 条样本，覆盖 next_probe_action / information_gain_judgment 等任务",
        "可直接喂给 Qwen3.5-4B/9B 与 Gemma-4-E4B 的 QLoRA 微调脚本（在 5090 服务器执行）",
    ])

    _add_game_matrix_slide(prs)

    # ---- Section: Roadmap ----
    _add_section_slide(prs, "09", "路线图")

    _add_harness_slide(prs)

    _add_windows_batch_slide(prs)

    _add_roadmap_slide(prs)

    # ---- Section: Next ----
    _add_section_slide(prs, "10", "下一步")

    _add_next_steps_slide(prs)

    # ---- Thanks ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), _W, _H)
    _set_fill(bg, _C_PRIMARY)
    # Decorative translucent shapes
    for idx in range(3):
        r = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10 + idx * 0.7), Inches(1 + idx * 0.5), Inches(1.5), Inches(1.5))
        _set_fill(r, _C_SECONDARY)
        r.fill.fore_color.brightness = 0.3
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(12.3), Inches(1.4))
    _set_text_style(tb.text_frame.paragraphs[0], "谢谢，欢迎讨论", Pt(48), bold=True, color=_C_LIGHT, align=PP_ALIGN.CENTER)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.67), Inches(4.3), Inches(2), Inches(0.08))
    _set_fill(accent, _C_ACCENT)

    # Add footers with page numbers
    total = len(prs.slides)
    for idx, slide in enumerate(prs.slides, start=1):
        if idx == 1:
            continue
        _add_footer(slide, idx, total)

    PPTX_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPTX_PATH))
    print(f"PPTX -> {PPTX_PATH}")


def main() -> None:
    md = REPORT_MD.read_text(encoding="utf-8")
    body = md_to_latex(md)
    write_tex(body)
    compile_tex()
    write_pptx()




def _add_windows_batch_slide(prs) -> None:
    """Slide: Windows Edge real-GPU rendering + batch results."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "Windows Edge 真 GPU 渲染：批量实验打通")
    _add_left_bar(slide, _C_SECONDARY)
    _add_content_bg(slide)

    _card(slide, 0.65, 1.45, 6.0, 2.3, "为什么必须换渲染环境", [
        "Cocos 3.8.5 游戏（kingshot/whiteout）要求完整 WebGL2",
        "WSL 软渲染 / D3D12 / Vulkan 直通均缺扩展 → Error 16405 白屏",
        "Windows Edge headless + D3D11 真 GPU：WebGL2 完整",
        "WSL 连接 Edge CDP（--remote-debugging-port=9222 --mute-audio）",
    ], _C_ACCENT)

    _card(slide, 6.95, 1.45, 6.0, 2.3, "工程改动", [
        "harness 支持 PLAYABLE_BROWSER_CDP 连接外部浏览器",
        "startup_timeout_ms 提到 45s（Unity/Cocos 加载 10-15s 出 canvas）",
        "声音关闭：Edge --mute-audio + harness mute() 双保险",
        "WSL adapter(9100) ↔ Windows harness 双向可达",
    ], _C_SUCCESS)

    _add_table_slide_raw(slide, 0.65, 3.95, 12.3, 2.6,
        ["游戏", "引擎", "terminal", "steps", "gameplay", "plans", "actions"],
        [
            ["tiles-survive-5ec610abcdff", "Unity", "OPERATOR_INTERRUPTED", "20", "9", "5", "14"],
            ["kingshot-c378f843e877", "Cocos", "OPERATOR_INTERRUPTED", "22", "6", "4", "14"],
            ["kingshot-94766e5d61dc", "Cocos", "BUDGET_EXHAUSTED", "240", "7", "4", "12"],
            ["whiteout-survival-87790941fd83", "Cocos", "OPERATOR_INTERRUPTED", "16", "5", "4", "9"],
            ["whiteout-survival-86420cdd2bbb", "Cocos", "RUNTIME_FAULT", "16", "7", "2", "10"],
            ["whiteout-survival-12ababda99c7", "Cocos", "OPERATOR_INTERRUPTED", "22", "10", "5", "17"],
            ["kingshot-9423402859e9", "Cocos", "OPERATOR_INTERRUPTED", "23", "10", "5", "17"],
            ["whiteout-survival-b64a7594f0c2", "Cocos", "OPERATOR_INTERRUPTED", "19", "8", "5", "14"],
            ["kingshot-2653755ff3a0", "Cocos", "OPERATOR_INTERRUPTED", "18", "13", "5", "14"],
            ["kingshot-0042aa74feb8", "Cocos", "OPERATOR_INTERRUPTED", "20", "10", "6", "15"],
            ["whiteout-survival-b47a4f071e9c", "Cocos", "OPERATOR_INTERRUPTED", "21", "10", "5", "15"],
            ["kingshot-6dd565baa02d", "Cocos", "OPERATOR_INTERRUPTED", "20", "8", "5", "13"],
            ["kingshot-33efef78d709", "Cocos", "OPERATOR_INTERRUPTED", "21", "14", "4", "16"],
            ["kingshot-f9a4fa1b6227", "Cocos", "OPERATOR_INTERRUPTED", "23", "10", "3", "16"],
            ["kingshot-830518bfdad4", "Cocos", "OPERATOR_INTERRUPTED", "26", "19", "4", "21"],
            ["whiteout-survival-c47908da8b11", "Cocos", "OPERATOR_INTERRUPTED", "15", "9", "3", "9"],
            ["kingshot-80b82b11529c", "Cocos", "OPERATOR_INTERRUPTED", "18", "7", "5", "13"],
            ["kingshot-ce59e2a9a7a3", "Cocos", "RUNTIME_FAULT", "2", "0", "0", "0"],
            ["kingshot-29345f023e9e", "Cocos", "BLOCKED_UNSAFE", "0", "0", "0", "0"],
        ])

    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(6.55), Inches(12.3), Inches(0.55))
    _set_fill(note, _C_LIGHT)
    _set_text_style(note.text_frame.paragraphs[0],
        "34 游戏累计 22 个产生真实 gameplay（21 个 Cocos）；最佳 kingshot-830518bfdad4：19 gameplay / 21 actions；batch10 首次 4/4 全胜。策略规范化器让 mimo-v2.5 大策略通过 harness 严格契约（planner_rejected=0）。",
        Pt(13), color=_C_DARK, align=PP_ALIGN.CENTER)

if __name__ == "__main__":
    main()
